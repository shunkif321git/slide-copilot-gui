from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def create_pptx_with_slides(filename: str, slides_data: list):
    prs = Presentation()

    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # タイトルとコンテンツのレイアウト

        # タイトル設定
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_info["title"]

        # 本文設定
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear() # 既存テキストをクリア

        for line in slide_info["content"].split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.LEFT

    prs.save(filename) 
